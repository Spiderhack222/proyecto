import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- CONFIGURACIÓN DE COLORES ---
BG_COLOR = RGBColor(13, 17, 23)      # #0D1117 (Dark)
NEON_CYAN = RGBColor(0, 255, 255)    # #00FFFF
TERM_GREEN = RGBColor(63, 185, 80)   # #3FB950
TEXT_WHITE = RGBColor(255, 255, 255) # #FFFFFF

def set_background(slide):
    """Aplica el fondo oscuro a la diapositiva."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_title(slide, text, font_size=40):
    """Añade título con estilo Neón."""
    title = slide.shapes.title
    title.text = text
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.name = "Impact" # O Arial Black
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = NEON_CYAN
        paragraph.alignment = PP_ALIGN.LEFT

def add_body(slide, text, font_size=18, color=TEXT_WHITE):
    """Añade cuerpo de texto estilo consola."""
    # Usamos el placeholder de cuerpo si existe, si no creamos uno
    if len(slide.placeholders) > 1:
        body = slide.placeholders[1]
    else:
        # Fallback si el layout es diferente
        body = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    
    tf = body.text_frame
    tf.text = text
    for paragraph in tf.paragraphs:
        paragraph.font.name = "Consolas"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color

def create_presentation():
    prs = Presentation()

    # --- SLIDE 1: PORTADA ---
    slide = prs.slides.add_slide(prs.slide_layouts[0]) # Title Slide
    set_background(slide)
    
    title = slide.shapes.title
    title.text = "BENCHMARKING DE ALGORITMOS\nDE ORDENAMIENTO EN C"
    for p in title.text_frame.paragraphs:
        p.font.name = "Impact"
        p.font.size = Pt(44)
        p.font.color.rgb = NEON_CYAN
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Análisis de Complejidad, Rendimiento y Estabilidad\n\n> Facultad de Ingeniería Eléctrica | UMSNH\n> Team: Iván Calderón | Diego Chávez | León Guzmán\n> Status: Compiled_Successfully"
    for p in subtitle.text_frame.paragraphs:
        p.font.name = "Consolas"
        p.font.color.rgb = TERM_GREEN
        p.font.size = Pt(16)

    # --- SLIDE 2: OBJETIVO ---
    slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
    set_background(slide)
    add_title(slide, "~/proyecto/objetivo.sh")
    add_body(slide, "$ cat description.txt\n\n1. Diseñar un banco de pruebas (Benchmark) en lenguaje C.\n2. Medir métricas de rendimiento real:\n   - Tiempo de ejecución (ms)\n   - Número de Comparaciones\n   - Número de Movimientos de memoria\n3. Contrastar la teoría O(n) vs. la práctica.\n4. Evaluar 4 algoritmos bajo estrés.")

    # --- SLIDE 3: LOS CONTENDIENTES ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_background(slide)
    add_title(slide, "LOS CONTENDIENTES (ALGORITMOS)")
    
    content = """
    1. INSERTION SORT [Rojo]
       - Complejidad: O(n^2)
       - Tipo: Elemental
    
    2. MERGE SORT [Azul]
       - Complejidad: O(n log n)
       - Tipo: Divide y Vencerás
    
    3. COUNTING SORT [Verde]
       - Complejidad: O(n)
       - Tipo: No Comparativo (Enteros)
    
    4. INTROSORT [Morado]
       - Complejidad: Híbrido
       - Nota: Quicksort + Heapsort + Insertion
    """
    add_body(slide, content)

    # --- SLIDE 4: SPECS ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_background(slide)
    add_title(slide, "SYSTEM_SPECS & DATASETS")
    add_body(slide, "HARDWARE/SW:\n- CPU: AMD Ryzen 3 3250U / Ryzen 5 3500U\n- OS: Ubuntu 22.04 / 24.04 LTS\n- Compilador: GCC\n\nDATASETS:\n- Tamaños: 100 hasta 7,500 elementos\n- Distribuciones: Uniforme, Ordenado, Reverso, Casi Ordenado, Duplicados")

    # --- SLIDE 5: RESULTADOS UNIFORME ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_background(slide)
    add_title(slide, "PRUEBA: ALEATORIO (UNIFORME)")
    add_body(slide, "[INSERTAR GRÁFICA 'TIEMPO - UNIFORME' AQUÍ]\n\nNOTAS:\n- 🔴 Insertion Sort: Se dispara exponencialmente O(n^2).\n- 🟢 Counting Sort: Línea plana. El más rápido O(n).\n- 🔵 Merge/Intro: Crecimiento logarítmico estable.")

    # --- SLIDE 6: RESULTADOS REVERSO ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_background(slide)
    add_title(slide, "PRUEBA DE ESTRÉS: ORDEN INVERSO")
    add_body(slide, "[INSERTAR GRÁFICA 'TIEMPO - REVERSO' AQUÍ]\n\nNOTAS:\n- ⚠️ El Reto: Datos totalmente invertidos.\n- 🛡️ Introsort: NO falla. Activa Heapsort.\n- ❌ Insertion Sort: Rendimiento crítico.")

    # --- SLIDE 7: RESULTADOS CASI ORDENADO ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_background(slide)
    add_title(slide, "PRUEBA: CASI ORDENADO")
    add_body(slide, "[INSERTAR GRÁFICA 'TIEMPO - CASI ORDENADO' AQUÍ]\n\n🏆 GANADOR: Insertion Sort\n- ¿Por qué? Mejor caso O(n).\n- Merge e Introsort son más lentos por overhead.")

    # --- SLIDE 8: COUNTING SORT ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_background(slide)
    add_title(slide, "EL ESPECIALISTA: COUNTING SORT")
    add_body(slide, "- Velocidad: Superior a todos O(n).\n- Comparaciones: 0 (Cero). No compara elementos.\n- Desventaja: Consume mucha RAM si el rango (k) es grande.\n- Uso: Exclusivo para enteros.")

    # --- SLIDE 9: ESTABILIDAD ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_background(slide)
    add_title(slide, "VERIFICACIÓN DE ESTABILIDAD")
    add_body(slide, "Insertion Sort: ✅ SI (Mantiene orden relativo)\nMerge Sort:     ✅ SI (Seguro para BD)\nCounting Sort:  ✅ SI (Por diseño)\nIntrosort:      ⚠️ NO (Intercambios rompen orden)")

    # --- SLIDE 10: CONCLUSIONES ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_background(slide)
    add_title(slide, "// CONCLUSIONS.C")
    add_body(slide, "/*\n 1. No existe el algoritmo perfecto. El contexto define al ganador.\n\n 2. Insertion Sort: Pésimo en aleatorio, INCREÍBLE en casi ordenado.\n\n 3. Introsort: La opción más segura (General Purpose).\n\n 4. Counting Sort: Velocidad absurda, limitado a enteros.\n*/")

    # --- SLIDE 11: RECOMENDACIONES ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_background(slide)
    add_title(slide, "IF (DATA == ?)")
    add_body(slide, "1. ¿Enteros y rango corto? -> SI -> Counting Sort\n\n2. ¿Casi ordenados?        -> SI -> Insertion Sort\n\n3. ¿Necesitas estabilidad? -> SI -> Merge Sort\n\n4. ¿Caso General?          -> ELSE -> Introsort")

    # --- SLIDE 12: CIERRE ---
    slide = prs.slides.add_slide(prs.slide_layouts[0]) # Title Slide
    set_background(slide)
    title = slide.shapes.title
    title.text = "GRACIAS_"
    for p in title.text_frame.paragraphs:
        p.font.name = "Consolas"
        p.font.size = Pt(60)
        p.font.color.rgb = NEON_CYAN
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Repositorio: [Insertar Link]\nFacultad de Ingeniería Eléctrica - UMSNH\n2024/2025"
    for p in subtitle.text_frame.paragraphs:
        p.font.color.rgb = TEXT_WHITE

    # Guardar
    prs.save('Benchmarking_Algoritmos_Final.pptx')
    print("¡Presentación generada con éxito: Benchmarking_Algoritmos_Final.pptx!")

if __name__ == "__main__":
    create_presentation()