#!/usr/bin/env python3
"""Convierte documentacion.html en documentacion.pptx.
- Extrae cada <div class="slide"> y crea una diapositiva.
- Inserta título, texto de paneles y todas las imágenes encontradas en la slide.
"""
import os
from pathlib import Path
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import pygments
from pygments import lexers
from pygments.formatters import ImageFormatter
from pygments import highlight
import tempfile
from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
HTML_PATH = ROOT / "documentacion.html"
OUT_PPTX = ROOT / "documentacion.pptx"

def text_from_panel(panel):
    # Extrae textos de párrafos, listas y código dentro del panel
    parts = []
    for p in panel.find_all(['h3','p']):
        txt = p.get_text(separator=' ', strip=True)
        if txt:
            parts.append(txt)
    # lists
    for ul in panel.find_all('ul'):
        for li in ul.find_all('li'):
            parts.append('- ' + li.get_text(separator=' ', strip=True))
    # Do not inline code blocks here; they will be rendered as images
    return '\n'.join(parts)


def add_title(slide, title_text, prs):
    left = Inches(0.4)
    top = Inches(0.25)
    width = prs.slide_width - Inches(0.8)
    height = Inches(0.9)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.name = 'Orbitron'
    p.font.color.rgb = RGBColor(0x00, 0xFF, 0xFF)  # neon cyan


def add_body(slide, body_text):
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9.5)
    height = Inches(4.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_bottom = 0
    tf.margin_top = 0
    tf.margin_left = 0
    tf.margin_right = 0
    lines = body_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = line
        else:
            p = tf.add_paragraph()
            # if starts with '- ' make it a bullet
            if line.startswith('- '):
                p.text = line[2:]
                p.level = 1
                p.font.size = Pt(14)
            elif line.startswith('CODE:'):
                # code block: add as separate paragraph with monospace
                code_block = line[len('CODE:'):]
                for code_line in code_block.split('\n'):
                    pc = tf.add_paragraph()
                    pc.text = code_line
                    pc.font.name = 'Courier New'
                    pc.font.size = Pt(9)
            else:
                p.text = line

    # style for all paragraphs
    for p in tf.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0xE6, 0xED, 0xF3)


def add_images(slide, img_tags, prs):
    if not img_tags:
        return
    # place images in a row below body; compute widths
    imgs = [tag.get('src') for tag in img_tags if tag.get('src')]
    imgs = [ROOT / src for src in imgs]
    imgs = [p for p in imgs if p.exists()]
    if not imgs:
        return
    margin_left = Inches(0.5)
    margin_right = Inches(0.5)
    available_width = prs.slide_width - margin_left - margin_right
    n = len(imgs)
    max_width = available_width / n
    top = Inches(5.5)
    height = Inches(1.6)
    for i, img in enumerate(imgs):
        left = margin_left + i * max_width
        try:
            slide.shapes.add_picture(str(img), left, top, width=max_width - Inches(0.15))
        except Exception as e:
            print(f"No se pudo insertar imagen {img}: {e}")


def main():
    if not HTML_PATH.exists():
        print(f"No se encontró {HTML_PATH}")
        return
    html = HTML_PATH.read_text(encoding='utf-8')
    soup = BeautifulSoup(html, 'lxml')

    prs = Presentation()
    # set widescreen approximate 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = soup.find_all('div', class_='slide')
    for idx, s in enumerate(slides, start=1):
        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)
        # background color from CSS #0D1117
        try:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(0x0D, 0x11, 0x17)
        except Exception:
            pass
        # title: try .title
        title_el = s.select_one('.title')
        header_h1 = s.find(['h1','h2'])
        if title_el:
            title_text = title_el.get_text(' ', strip=True)
        elif header_h1:
            title_text = header_h1.get_text(' ', strip=True)
        else:
            title_text = f"Diapositiva {idx}"
        add_title(slide, title_text, prs)

        # collect panel texts
        panels = s.find_all('div', class_='panel')
        body_parts = []
        # collect code images to render
        code_image_paths = []
        for p in panels:
            # detect code-body blocks and render them to images
            for j, cb in enumerate(p.select('.code-body')):
                code_text = cb.get_text('\n', strip=True)
                if code_text:
                    # choose lexer (C) or guess
                    try:
                        lexer = lexers.CLexer()
                    except Exception:
                        lexer = lexers.get_lexer_by_name('c')
                    formatter = ImageFormatter(font_name='DejaVu Sans Mono', font_size=14, line_numbers=False, style='monokai')
                    try:
                        img_bytes = highlight(code_text, lexer, formatter)
                        tmpf = Path(tempfile.gettempdir()) / f"slide{idx}_code{j}.png"
                        with open(tmpf, 'wb') as fh:
                            fh.write(img_bytes)
                        code_image_paths.append(tmpf)
                    except Exception as e:
                        # fallback: skip rendering
                        print(f"No se pudo renderizar código en imagen: {e}")
            t = text_from_panel(p)
            if t:
                body_parts.append(t)
        body_text = '\n\n'.join(body_parts)[:15000]
        if body_text:
            add_body(slide, body_text)

        # images from HTML
        img_tags = s.find_all('img')
        add_images(slide, img_tags, prs)

        # add rendered code images
        if code_image_paths:
            # place below any existing images
            add_images(slide, [BeautifulSoup(f'<img src="{p}"></img>', 'lxml').img for p in code_image_paths], prs)

        # add footer / page number
        try:
            footer = slide.shapes.add_textbox(prs.slide_width - Inches(2.2), prs.slide_height - Inches(0.5), Inches(2), Inches(0.4))
            ft = footer.text_frame
            ft.clear()
            pnum = ft.paragraphs[0]
            pnum.text = f"Página {idx}"
            pnum.font.size = Pt(10)
            pnum.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
            pnum.font.name = 'Fira Code'
        except Exception:
            pass

    prs.save(OUT_PPTX)
    print(f"Guardado: {OUT_PPTX}")

if __name__ == '__main__':
    main()
